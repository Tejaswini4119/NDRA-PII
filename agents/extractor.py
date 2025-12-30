
import hashlib
import os
import mimetypes
import uuid
import re
import json
import yaml
import zipfile
import tarfile
from typing import List, Dict, Any, Optional
from datetime import datetime
from pathlib import Path
import logging

# External libs
from pypdf import PdfReader
from docx import Document as DocxDocument
import pandas as pd
from pptx import Presentation
from bs4 import BeautifulSoup
from PIL import Image
try:
    import extract_msg
except ImportError:
    extract_msg = None

# Internal
from agents.base import NDRAAgent
from schemas.core_models import DocumentMetadata, SemanticChunk, RawChunk

class ExtractorAgent(NDRAAgent):
    """
    Phase 2 (Expanded): Multi-Model Ingestion & Feature Extraction.
    Supports: 
    - Docs: PDF, DOCX, PPTX
    - Data: TXT, CSV, JSON, XML, YAML, LOG, SQL
    - Web: HTML
    - Email: EML, MSG
    - Images: PNG, JPG (Metadata only for now, OCR stub)
    - Archives: ZIP, TAR (Listing)
    - Quarantine: Safe fallback for malformed/unsupported.
    """
    
    def __init__(self, quarantine_dir="quarantine"):
        super().__init__("ExtractorAgent")
        self.quarantine_dir = Path(quarantine_dir)
        self.quarantine_dir.mkdir(exist_ok=True)
        
        # Dispatch Table (Mime/Ext -> Handler)
        self.handlers = {
            # Documents
            "application/pdf": self._read_pdf,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": self._read_docx,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": self._read_pptx,
            
            # Spreadsheet / Tabular
            "text/csv": self._read_csv,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": self._read_excel,
            "application/vnd.ms-excel": self._read_excel,
            "application/json": self._read_json,
            "application/xml": self._read_xml,
            "text/xml": self._read_xml,
            "application/x-yaml": self._read_yaml,
            "text/yaml": self._read_yaml,
            
            # Text / Code
            "text/plain": self._read_txt,
            "text/markdown": self._read_txt,
            "text/html": self._read_html,
            
            # Image
            "image/jpeg": self._read_image,
            "image/png": self._read_image,
            "image/webp": self._read_image,
            
            # Email
            "message/rfc822": self._read_eml,
            "application/vnd.ms-outlook": self._read_msg,
            
            # Archive
            "application/zip": self._read_archive_zip,
            "application/x-tar": self._read_archive_tar,
            "application/gzip": self._read_archive_tar,
        }
        
        # Extensions fallback if mime guessing fails or is generic
        self.ext_map = {
            ".md": "text/markdown",
            ".log": "text/plain",
            ".sql": "text/plain",
            ".yml": "text/yaml",
            ".yaml": "text/yaml",
            ".py": "text/plain",
            ".msg": "application/vnd.ms-outlook",
            ".eml": "message/rfc822",
            ".json": "application/json",
            ".parquet": "application/octet-stream" # Needs specific handler
        }

    def process(self, file_path: str, context: Dict[str, Any] = None) -> List[SemanticChunk]:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        # 1. Integrity & Type
        file_hash = self._compute_sha256(path)
        mime_type = self._detect_mime(path)
        
        doc_meta = DocumentMetadata(
            filename=path.name,
            file_size_bytes=path.stat().st_size,
            mime_type=mime_type,
            sha256_hash=file_hash,
            source_channel="file_system"
        )
        self.log_event("DOCUMENT_RECEIVED", doc_meta.dict())

        # 2. Select Handler or Quarantine
        handler = self.handlers.get(mime_type)
        if not handler:
            self._quarantine_file(path, "Unsupported MIME type")
            return []

        # 3. Extract
        try:
            raw_chunks = handler(path)
            if not raw_chunks:
                 self.logger.warning(f"No text extracted from {path.name}")
                 return []
                 
            # 4. Semantic Chunking
            semantic_chunks = self._chunk_text(raw_chunks, doc_meta)
            
            self.log_event("EXTRACTION_COMPLETE", {
                "file": path.name,
                "mime": mime_type,
                "chunks": len(semantic_chunks)
            })
            return semantic_chunks
            
        except Exception as e:
            self.logger.error(f"Extraction failed for {path.name}: {e}")
            self._quarantine_file(path, f"Extraction Error: {str(e)}")
            return []

    # --- Quarantine ---
    def _quarantine_file(self, original_path: Path, reason: str):
        self.logger.warning(f"Quarantining {original_path.name}: {reason}")
        # In real system, move file. Here just log logic.
        dest = self.quarantine_dir / f"{original_path.name}_{uuid.uuid4().hex[:6]}"
        # shutil.copy(original_path, dest) # Uncomment to actually copy
        self.log_event("FILE_QUARANTINED", {"file": original_path.name, "reason": reason})

    # --- Handlers ---
    
    def _read_pdf(self, path: Path) -> List[Dict]:
        reader = PdfReader(path)
        return [{"text": p.extract_text() or "", "page": i+1} for i, p in enumerate(reader.pages)]

    def _read_docx(self, path: Path) -> List[Dict]:
        doc = DocxDocument(path)
        text = "\n".join([p.text for p in doc.paragraphs])
        return [{"text": text, "page": 1}] # Structurally one unit

    def _read_pptx(self, path: Path) -> List[Dict]:
        prs = Presentation(path)
        chunks = []
        for i, slide in enumerate(prs.slides):
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
            chunks.append({"text": "\n".join(text), "page": i+1})
        return chunks

    def _read_txt(self, path: Path) -> List[Dict]:
        # Try utf-8 then latin-1
        try:
            with open(path, "r", encoding="utf-8") as f: return [{"text": f.read(), "page": 1}]
        except UnicodeDecodeError:
             with open(path, "r", encoding="latin-1") as f: return [{"text": f.read(), "page": 1}]

    def _read_csv(self, path: Path) -> List[Dict]:
        df = pd.read_csv(path)
        text = df.to_csv(index=False) # Convert back to string for text analysis (or handle as structured)
        return [{"text": text, "page": 1}]

    def _read_excel(self, path: Path) -> List[Dict]:
        sheets = pd.read_excel(path, sheet_name=None)
        chunks = []
        for name, df in sheets.items():
            text = f"Sheet: {name}\n" + df.to_string()
            chunks.append({"text": text, "page": 1})
        return chunks

    def _read_json(self, path: Path) -> List[Dict]:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return [{"text": json.dumps(data, indent=2), "page": 1}]
    
    def _read_yaml(self, path: Path) -> List[Dict]:
        with open(path, "r", encoding="utf-8") as f:
            data = yaml.safe_load(f)
        return [{"text": yaml.dump(data), "page": 1}]
        
    def _read_xml(self, path: Path) -> List[Dict]:
         with open(path, "r", encoding="utf-8") as f:
             soup = BeautifulSoup(f, "xml")
             return [{"text": soup.get_text(separator="\n"), "page": 1}]

    def _read_html(self, path: Path) -> List[Dict]:
        with open(path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f, "html.parser")
            return [{"text": soup.get_text(separator="\n"), "page": 1}]

    def _read_image(self, path: Path) -> List[Dict]:
        # Phase 2 stub for OCR or metadata
        try:
            img = Image.open(path)
            meta_text = f"[IMAGE FILE: {path.name}]\nFormat: {img.format}\nSize: {img.size}\nMode: {img.mode}"
            # TODO: Integrate PaddleOCR here if requested
            return [{"text": meta_text, "page": 1}]
        except Exception:
            return []

    def _read_eml(self, path: Path) -> List[Dict]:
        # Basic parsing of EML
        import email
        with open(path, "rb") as f:
            msg = email.message_from_binary_file(f)
        
        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                content_type = part.get_content_type()
                if content_type == "text/plain":
                    body += part.get_payload(decode=True).decode()
        else:
            body = msg.get_payload(decode=True).decode()
            
        full_text = f"Subject: {msg['subject']}\nFrom: {msg['from']}\nTo: {msg['to']}\n\n{body}"
        return [{"text": full_text, "page": 1}]

    def _read_msg(self, path: Path) -> List[Dict]:
        if extract_msg:
            msg = extract_msg.Message(path)
            return [{"text": msg.body, "page": 1}]
        return [{"text": "[MSG Support Missing - Install extract-msg]", "page": 1}]

    def _read_archive_zip(self, path: Path) -> List[Dict]:
        # Returns list of files contained (Metadata), usually don't extract recursively in 'chunk' mode
        # unless we treat archive as a folder. For now, list contents.
        with zipfile.ZipFile(path, 'r') as z:
            names = z.namelist()
        return [{"text": f"Archive Contents ({path.name}):\n" + "\n".join(names), "page": 1}]

    def _read_archive_tar(self, path: Path) -> List[Dict]:
         with tarfile.open(path, 'r') as t:
            names = t.getnames()
         return [{"text": f"Archive Contents ({path.name}):\n" + "\n".join(names), "page": 1}]

    # --- Utils ---
    def _compute_sha256(self, file_path: Path) -> str:
        sha256_hash = hashlib.sha256()
        with open(file_path, "rb") as f:
            for byte_block in iter(lambda: f.read(4096), b""):
                sha256_hash.update(byte_block)
        return sha256_hash.hexdigest()

    def _detect_mime(self, file_path: Path) -> str:
        # Check explicit ext map first
        ext = file_path.suffix.lower()
        if ext in self.ext_map:
            return self.ext_map[ext]
            
        mime, _ = mimetypes.guess_type(file_path)
        return mime or "application/octet-stream"

    def _chunk_text(self, raw_pages: List[Dict], meta: DocumentMetadata) -> List[SemanticChunk]:
        """
        Sliding window semantic chunking (Same as Phase 2).
        """
        CHUNK_SIZE = 800 # Increased for Data files
        OVERLAP = 100
        
        chunks = []
        doc_id = meta.sha256_hash
        
        for entry in raw_pages:
            page_text = str(entry["text"]) # Ensure string
            page_num = entry["page"]
            
            # Normalize
            page_text = re.sub(r'\s+', ' ', page_text).strip()
            
            start = 0
            while start < len(page_text):
                end = min(start + CHUNK_SIZE, len(page_text))
                
                # Check sentence boundary
                if end < len(page_text):
                    last_period = page_text.rfind('.', start, end)
                    if last_period != -1 and (end - last_period) < 100:
                        end = last_period + 1
                        
                chunk_text = page_text[start:end]
                
                chunk = SemanticChunk(
                   document_id=doc_id,
                   processed_text=chunk_text,
                   original_text=chunk_text,
                   page_number=page_num,
                   token_span=(start, end),
                   section_label="content" 
                )
                chunks.append(chunk)
                
                start += (CHUNK_SIZE - OVERLAP)
                
        return chunks
